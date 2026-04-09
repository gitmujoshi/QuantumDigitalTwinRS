#!/usr/bin/env python3
"""Generate TwinSentry-RS overview deck (requires: pip install python-pptx)."""

from __future__ import annotations

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError as e:
    raise SystemExit("Install python-pptx: pip install python-pptx") from e


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets(prs: Presentation, title: str, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "TwinSentry-RS",
        "Digital Twin Control Plane — Training Overview\nQuantum & Software Tracks",
    )

    add_bullets(
        prs,
        "Learning objectives",
        [
            "Align quantum and software roles on one architecture",
            "Know what the twin simulates — and what it does not",
            "Trace intent → BAML → Rust → Langfuse end-to-end",
            "Build the Rust core and Python extension locally",
        ],
    )

    add_bullets(
        prs,
        "What is TwinSentry-RS?",
        [
            "Bridges AI intent and low-level pulse execution",
            "Policy-validated pulses (BAML) + deterministic simulation (Rust)",
            "Observable runs via self-hosted Langfuse on GCP/VPC",
            "Rust simulation core has no network — queue-only ingress",
        ],
    )

    add_bullets(
        prs,
        "Three pillars",
        [
            "Safety: BAML schema + asserts / checks (PEP-style enforcement)",
            "Fidelity: 2-qubit TDSE, RK4, SIMD-friendly ℂ⁴ ops, SPSC updates",
            "Auditability: full trace ID from intent to fidelity score",
        ],
    )

    add_bullets(
        prs,
        "Architecture (conceptual)",
        [
            "Control plane: Python, BAML, Langfuse, future Streamlit UI",
            "Data plane: Rust crate twin_sentry — state, Hamiltonian, RK4",
            "Bridge: PyO3 extension twin_sentry._native",
            "Observability: docker-compose Langfuse stack (Postgres, ClickHouse, …)",
        ],
    )

    add_bullets(
        prs,
        "Quantum track — State & evolution",
        [
            "2-qubit ket in ℂ⁴: |00⟩ … |11⟩ (computational basis)",
            "dψ/dt = −i H(t) ψ — integrated with RK4 (see src/rk4.rs)",
            "Hamiltonian: Kronecker-built 4×4 H; drive + Zeeman-style splits",
            "fidelity_ground ≈ population of |00⟩ — simple proxy, not full tomography",
        ],
    )

    add_bullets(
        prs,
        "Quantum track — Pulses & policy",
        [
            "BAML: QuantumPulse (amplitude, frequency, duration, gate_type)",
            "Optional NoiseProfile: T2 / thermal as capped relative scales",
            "Asserts: hard-stop invalid outputs; checks: inspectable policy bands",
            "Rust PulseCommand: maps to Hamiltonian parameters for the twin",
        ],
    )

    add_bullets(
        prs,
        "Software track — Repository map",
        [
            "baml_src/*.baml — contracts & Gemini client",
            "baml_client/ — generated Python (baml-cli generate)",
            "src/*.rs — physics, SIMD ops, SPSC, twin, bridge (PyO3)",
            "python/twin_sentry/ — package + controller.py",
            "docker-compose.yaml, cloudbuild.yaml — Langfuse & CI",
        ],
    )

    add_bullets(
        prs,
        "Software track — Build & run",
        [
            "Rust: cargo test --all-features (needs python3-dev on Linux for PyO3)",
            "Python: maturin develop --features python",
            "Env: GOOGLE_API_KEY (BAML), LANGFUSE_* (tracing)",
            "Controller: run_twin_pipeline() — spans + fidelity score",
        ],
    )

    add_bullets(
        prs,
        "Joint — End-to-end flow",
        [
            "1. NL intent → 2. BAML ParsePulseFromIntent",
            "3. Map to PulseCommand + noise metadata to Langfuse",
            "4. TwinEngine.step / drain / renormalize",
            "5. Score fidelity_ground on trace",
        ],
    )

    add_bullets(
        prs,
        "Hands-on ideas",
        [
            "Quantum: vary dt / steps; watch fidelity stability",
            "Software: import TwinEngine; push PulseCommand via pulse_queue",
            "Joint: run Langfuse locally + one traced pipeline",
        ],
    )

    add_bullets(
        prs,
        "Documentation & next steps",
        [
            "Full narrative: docs/training/TwinSentry-Training-Guide.md",
            "Regenerate this deck: python docs/training/build_slides.py",
            "Next milestones: Streamlit UI, richer noise, benchmarks vs QuTiP",
        ],
    )

    out = Path(__file__).resolve().parent / "TwinSentry-Training-Overview.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
